from django.shortcuts import render,redirect,HttpResponse
from django.contrib.auth.decorators import login_required

from .models import *
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt



@login_required
def next(request,pk):
    obj = pitchDeck.objects.get(pitchDeckId=pk)
    prs = Presentation()
    
    ##slide1 titleslide
    main = titleSlide.objects.get(pitchDeckId=pk)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = main.name 
    subtitle.text = main.titleline


    ##slide2 problemTryingToSolve
    '''problem = problemTryingToSolve.objects.get(pitchDeckId=pk)
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Problem'

    tf = body_shape.text_frame
    tf.text = problem.bullet1

    p = tf.add_paragraph()
    p.text = problem.bullet2
    p.level = 0

    p = tf.add_paragraph()
    p.text = problem.bullet3
    p.level = 0

    p = tf.add_paragraph()
    p.text = problem.bullet4
    p.level = 0

    p = tf.add_paragraph()
    p.text = problem.bullet5
    p.level = 0'''


    ##slide3 vision
    Vision = vision.objects.get(pitchDeckId=pk)
    
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Vision'

    tf = body_shape.text_frame
    tf.text = Vision.bullet1

    p = tf.add_paragraph()
    p.text = Vision.bullet2
    p.level = 0

    p = tf.add_paragraph()
    p.text = Vision.bullet3
    p.level = 0

    p = tf.add_paragraph()
    p.text = Vision.bullet4
    p.level = 0

    p = tf.add_paragraph()
    p.text = Vision.bullet5
    p.level = 0

    ##slide4 uniqueValueProposition##
    USP = uniqueValueProposition.objects.get(pitchDeckId=pk)
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = USP.usp

    ##autoshapes here
    width = Inches(2.2)
    height = Inches(2.2)

    #1
    left = Inches(3.3)  
    top = Inches(2.00)
    shape1 = shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    shape1.rotation = 90

    #2
    left = Inches(2.1)  
    top = Inches(3.8)
    shape2 = shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    shape2.rotation = 90

    #3
    left = Inches(4.5)  # 0.93" centers this overall set of shapes
    top = Inches(3.8)
    shape3 = shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    shape3.rotation = 90

    #textbox

    width_tb = Inches(2.17)
    height_tb = Inches(0.7)
    #1
    left_tb =Inches(3.31)
    top_tb = Inches(2.74)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    tf1.text = USP.bullet1
    #2
    left_tb = Inches(2.15)
    top_tb = Inches(4.5)
    txBox2 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf2 = txBox2.text_frame
    tf2.text = USP.bullet2

    #3
    left_tb = Inches(4.53)
    top_tb = Inches(4.5)
    txBox3 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf3 = txBox3.text_frame
    tf3.text = USP.bullet3

    ##slide5 underlyingMagic
    
    ##slide6 milestones
    

    ##slide7 totalAddressableMarket
    TAM = totalAddressableMarket.objects.get(pitchDeckId=pk)
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Total Addressable Market'

    ##autoshapes here

    #1
    width = Inches(5.2)
    height = Inches(5.2)
    left = Inches(2.3)  
    top = Inches(1.8)
    shape1 = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape1.rotation = 270
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #2

    width = Inches(4.0)
    height = Inches(4.0)
    left = Inches(2.9)  
    top = Inches(3.0)
    shape2 = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill = shape2.fill
    fill.solid()
    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #3

    width = Inches(2.5)
    height = Inches(2.5)
    left = Inches(3.6)  
    top = Inches(4.5)
    shape3 = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape3.rotation = 90
    fill = shape3.fill
    fill.solid()
    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    fill.fore_color.brightness = -0.25


    #textbox

    #1
    width_tb = Inches(3.20)
    height_tb = Inches(0.80)
    left_tb =Inches(3.4)
    top_tb = Inches(2.38)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    tf1.text = TAM.totalAvailableMarket
    tf1.alignment = PP_ALIGN.CENTER

    #2
    width_tb = Inches(3.5)
    height_tb = Inches(0.9)
    left_tb = Inches(3.2)
    top_tb = Inches(3.70)
    txBox2 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf2 = txBox2.text_frame
    tf2.text = TAM.servedAvailableMarket
    tf2.alignment = PP_ALIGN.CENTER


    #3
    width_tb = Inches(2.1)
    height_tb = Inches(0.9)
    left_tb = Inches(3.9)
    top_tb = Inches(5.3)
    txBox3 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf3 = txBox3.text_frame
    tf3.text = TAM.targetMarket


    ##slide8 businessModal
    bM = businessModal.objects.get(pitchDeckId = pk)
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes


    ##autoshapes here

    width = Inches(3.4)
    height = Inches(2.1)

    #1
    left = Inches(0.3)  
    top = Inches(0.6)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #2
    left = Inches(6.2)  
    top = Inches(0.6)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #3
    left = Inches(0.3)  
    top = Inches(4.5)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #4
    left = Inches(6.2)  
    top = Inches(4.5)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1



    #list Textbox

    width_tb = Inches(3.4)
    height_tb = Inches(0.66)
    #1
    left_tb = Inches(0.3)
    top_tb = Inches(1.3)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    tf1.text = bM.customerAcquisitionStrategy
    tf1.alignment = PP_ALIGN.CENTER


    #2
    left_tb = Inches(6.2)
    top_tb = Inches(1.24)
    txBox2 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf2 = txBox2.text_frame
    tf2.text =  bM.serviceProvision 
    tf2.alignment = PP_ALIGN.CENTER

    #3
    left_tb = Inches(6.2)
    top_tb = Inches(5.24)
    txBox3 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf3 = txBox3.text_frame
    tf3.text = bM.modelType
    tf3.alignment = PP_ALIGN.CENTER


    #4
    left_tb = Inches(0.3)
    top_tb = Inches(5.24)
    txBox4 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf4 = txBox4.text_frame
    tf4.text = bM.monetizationStrategy
    tf4.alignment = PP_ALIGN.CENTER

    ##Title


    width_tb = Inches(6.3)
    height_tb = Inches(0.8)
    left_tb = Inches(2.0)
    top_tb = Inches(3.1)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    p = tf1.paragraphs[0]
    run = p.add_run()
    run.text = 'Business Model'
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(44)


    ##slide9 Competition
    C = Competition.objects.get(pitchDeckId=pk)
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Competitors'

    ##autoshapes here

    width = Inches(2.5)
    height = Inches(2.1)

    #1
    left = Inches(2.9)  
    top = Inches(2.1)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #2
    left = Inches(5.7)  
    top = Inches(2.1)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #3
    left = Inches(2.9)  
    top = Inches(4.5)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #4
    left = Inches(5.7)  
    top = Inches(4.5)
    shape1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape1.fill
    fill.gradient()
    fill.gradientangle = 90
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.ACCENT_1

    #textbox

    width_tb = Inches(1.6)
    height_tb = Inches(0.38)

    #1
    left_tb =Inches(1.1)
    top_tb = Inches(2.7)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    tf1.text = "Startup"
    tf1.alignment = PP_ALIGN.CENTER

    #2
    left_tb = Inches(1.1)
    top_tb = Inches(5.30)
    txBox2 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb) 
    tf2 = txBox2.text_frame
    tf2.text = "Incumbents"
    tf2.alignment = PP_ALIGN.CENTER


    #3
    left_tb = Inches(3.4)
    top_tb = Inches(1.42)
    txBox3 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf3 = txBox3.text_frame
    tf3.text = "Paid"
    tf3.alignment = PP_ALIGN.CENTER


    #4
    left_tb = Inches(6.10)
    top_tb = Inches(1.42)
    txBox4 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf4 = txBox4.text_frame
    tf4.text = "Free"
    tf4.alignment = PP_ALIGN.CENTER


    #list Textbox

    width_tb = Inches(1.2)
    height_tb = Inches(0.94)
    #1
    left_tb = Inches(3.4)
    top_tb = Inches(2.7)
    txBox1 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf1 = txBox1.text_frame
    tf1.text = C.startupsFree
    tf1.alignment = PP_ALIGN.CENTER


    #2
    left_tb = Inches(6.3)
    top_tb = Inches(2.7)
    txBox2 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf2 = txBox2.text_frame
    tf2.text = C.startupsPaid
    tf2.alignment = PP_ALIGN.CENTER


    #3
    left_tb = Inches(3.4)
    top_tb = Inches(4.96)
    txBox3 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf3 = txBox3.text_frame
    tf3.text = C.incumbentsFree 
    tf3.alignment = PP_ALIGN.CENTER


    #4
    left_tb = Inches(6.3)
    top_tb = Inches(4.95)
    txBox4 = slide.shapes.add_textbox(left_tb, top_tb, width_tb, height_tb)
    tf4 = txBox4.text_frame
    tf4.text = C.incumbentsPaid
    tf4.alignment = PP_ALIGN.CENTER



    ###slide10 competitiveAdvantage





    prs.save(obj.startupName+".pptx")
    return render(request,'generate.html',{'obj':obj})

